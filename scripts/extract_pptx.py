"""Extract slide text and speaker notes from every .pptx in 'Old artefacts AI Product Manager '.

Writes one markdown file per deck into scripts/_out/<module>.md with sections per slide,
preserving slide order, body text, and speaker notes. Used as the raw source we synthesise
the new Slides.html / Notes.md / Frameworks / Glossary / Pre-Read from.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Install requirements first: pip install -r scripts/requirements.txt", file=sys.stderr)
    sys.exit(1)

ROOT = Path(__file__).resolve().parent.parent
SOURCE = ROOT / "Old artefacts AI Product Manager "
OUT = ROOT / "scripts" / "_out"
OUT.mkdir(parents=True, exist_ok=True)


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes = slide.notes_slide.notes_text_frame
    if not notes:
        return ""
    return notes.text.strip()


def slugify(name: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_").lower()
    return s[:60]


def dump_deck(pptx_path: Path) -> Path:
    prs = Presentation(pptx_path)
    out = OUT / f"{slugify(pptx_path.stem)}.md"
    with out.open("w") as f:
        f.write(f"# {pptx_path.name}\n\n")
        for idx, slide in enumerate(prs.slides, start=1):
            body = slide_text(slide)
            notes = notes_text(slide)
            f.write(f"\n---\n\n## Slide {idx}\n\n")
            if body:
                f.write("### Body\n\n")
                f.write(body + "\n\n")
            if notes:
                f.write("### Speaker Notes\n\n")
                f.write(notes + "\n\n")
    return out


def main() -> None:
    if not SOURCE.exists():
        print(f"Source folder not found: {SOURCE}", file=sys.stderr)
        sys.exit(1)
    decks = sorted(SOURCE.rglob("*.pptx"))
    if not decks:
        print("No .pptx files found.", file=sys.stderr)
        sys.exit(1)
    for deck in decks:
        out = dump_deck(deck)
        print(f"wrote {out.relative_to(ROOT)}  ({deck.relative_to(ROOT)})")


if __name__ == "__main__":
    main()
